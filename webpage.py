from pptx import Presentation

# Define the slide title and bullet points
slide_title = 'Introduction'
bullet_points = [
    'Definition of game theory',
    'Examples of game theory applications (e.g., economics, politics, biology)',
    'Importance of auctions and negotiation in game theory'
]

# Create the slide and add the title and bullet points
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = slide_title
left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
for point in bullet_points:
    text_frame.add_paragraph().text = point

# Define the slide title and bullet points
slide_title = 'Game Theory Basics'
bullet_points = [
    'Definition of games and players',
    'Types of games (e.g., simultaneous vs. sequential, zero-sum vs. non-zero-sum)',
    'Nash equilibrium and its importance'
]

# Create the slide and add the title and bullet points
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = slide_title
left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
for point in bullet_points:
    text_frame.add_paragraph().text = point